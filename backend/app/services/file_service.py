import os
import uuid
import shutil
from typing import List, Tuple, Optional
from datetime import datetime
from pathlib import Path
import fitz  # PyMuPDF
from pptx import Presentation
from PIL import Image
import io
import asyncio

class FileService:
    def __init__(self):
        # Use data directory for all data storage (relative to backend folder)
        base_dir = Path.cwd()
        self.upload_dir = base_dir / "data" / "projects"
        self.temp_dir = base_dir / "data" / "temp"
        self.upload_dir.mkdir(parents=True, exist_ok=True)
        self.temp_dir.mkdir(parents=True, exist_ok=True)
        self.websocket_manager = None
        self.session_id = None
    
    def set_websocket_session(self, manager, session_id: str):
        """Set WebSocket manager and session ID for progress updates"""
        self.websocket_manager = manager
        self.session_id = session_id
    
    async def send_progress(self, icon: str, message: str):
        """Send progress update via WebSocket"""
        if self.websocket_manager and self.session_id:
            try:
                await self.websocket_manager.send_message(self.session_id, {
                    "type": "progress",
                    "icon": icon,
                    "message": message
                })
            except Exception as e:
                print(f"❌ Failed to send progress update: {e}")
        
    async def process_uploaded_file(self, file_path: str, filename: str) -> Tuple[str, List[str], List[str]]:
        """
        Process uploaded presentation file and extract slides
        Returns: (project_id, slide_image_paths, extracted_texts)
        """
        try:
            await self.send_progress("🔄", "Processing file...")
            
            # Get file size for progress info
            file_size = os.path.getsize(file_path)
            await self.send_progress("📊", f"File size: {file_size / (1024*1024):.2f} MB")
            
            # Generate unique project ID
            project_id = str(uuid.uuid4())
            project_dir = self.upload_dir / project_id
            project_dir.mkdir(exist_ok=True)
            
            await self.send_progress("💾", "Creating temporary file...")
            
            # Create subdirectories
            (project_dir / "slides").mkdir(exist_ok=True)
            (project_dir / "scripts").mkdir(exist_ok=True)
            (project_dir / "audio").mkdir(exist_ok=True)
            
            # Copy original file
            original_file_path = project_dir / filename
            shutil.copy2(file_path, original_file_path)
            
            await self.send_progress("✅", "Temporary file created")
            
            # Detect file type and process
            file_extension = Path(filename).suffix.lower()
            await self.send_progress("🔍", f"File extension: {file_extension}")
            
            if file_extension == '.pptx':
                slide_images, extracted_texts = await self._process_pptx(original_file_path, project_dir)
            elif file_extension == '.pdf':
                slide_images, extracted_texts = await self._process_pdf(original_file_path, project_dir)
            elif file_extension in ['.png', '.jpg', '.jpeg']:
                slide_images, extracted_texts = await self._process_image(original_file_path, project_dir)
            else:
                raise ValueError(f"Unsupported file type: {file_extension}")
            
            await self.send_progress("📝", "Saving project metadata...")
            
            # Create project metadata
            metadata = {
                "project_id": project_id,
                "original_filename": filename,
                "file_type": file_extension,
                "file_size": file_size,
                "created_at": datetime.now().isoformat(),
                "total_slides": len(slide_images),
                "slide_images": slide_images,
                "extracted_texts": extracted_texts
            }
            
            # Save metadata
            metadata_path = project_dir / "metadata.json"
            import json
            with open(metadata_path, 'w', encoding='utf-8') as f:
                json.dump(metadata, f, indent=2, ensure_ascii=False)
            
            await self.send_progress("🎉", f"Processing complete! Generated {len(slide_images)} slides")
            
            return project_id, slide_images, extracted_texts
            
        except Exception as e:
            await self.send_progress("❌", f"Error: {str(e)}")
            raise Exception(f"File processing failed: {str(e)}")
    
    async def _process_pptx(self, file_path: Path, project_dir: Path) -> Tuple[List[str], List[str]]:
        """
        Process PowerPoint file
        """
        try:
            # Extract text using python-pptx
            prs = Presentation(file_path)
            extracted_texts = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                extracted_texts.append(slide_text.strip())
            
            # Convert slides to images
            slide_images = await self._convert_pptx_to_images(file_path, project_dir)
            
            return slide_images, extracted_texts
            
        except Exception as e:
            raise Exception(f"PPTX processing failed: {str(e)}")
    
    async def _process_pdf(self, file_path: Path, project_dir: Path) -> Tuple[List[str], List[str]]:
        """
        Process PDF file
        """
        try:
            doc = fitz.open(file_path)
            slide_images = []
            extracted_texts = []
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # Extract text
                text = page.get_text()
                extracted_texts.append(text.strip())
                
                # Convert page to image
                mat = fitz.Matrix(2.0, 2.0)  # 2x zoom for better quality
                pix = page.get_pixmap(matrix=mat)
                img_data = pix.tobytes("png")
                
                # Save image
                image_filename = f"slide_{page_num + 1:03d}.png"
                image_path = project_dir / "slides" / image_filename
                
                with open(image_path, "wb") as f:
                    f.write(img_data)
                
                slide_images.append(str(image_path))
            
            doc.close()
            return slide_images, extracted_texts
            
        except Exception as e:
            raise Exception(f"PDF processing failed: {str(e)}")
    
    async def _process_image(self, file_path: Path, project_dir: Path) -> Tuple[List[str], List[str]]:
        """
        Process single image file
        """
        try:
            # Copy image to slides directory
            image_filename = "slide_001.png"
            image_path = project_dir / "slides" / image_filename
            
            # Convert and optimize image
            with Image.open(file_path) as img:
                # Convert to RGB if necessary
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                
                # Resize if too large
                if img.width > 1920 or img.height > 1080:
                    img.thumbnail((1920, 1080), Image.Resampling.LANCZOS)
                
                img.save(image_path, "PNG", optimize=True)
            
            return [str(image_path)], [""]  # No text extraction for single images
            
        except Exception as e:
            raise Exception(f"Image processing failed: {str(e)}")
    
    async def _convert_pptx_to_images(self, pptx_path: Path, project_dir: Path) -> List[str]:
        """
        Convert PPTX slides to high-quality images using LibreOffice + pdfplumber
        macOS optimized approach for best quality
        """
        try:
            slides_dir = project_dir / "slides"
            
            # Method 1: LibreOffice + pdfplumber (highest quality)
            slide_images = await self._convert_via_libreoffice_pdf(pptx_path, project_dir)
            if slide_images:
                print(f"✅ LibreOffice + pdfplumber conversion successful: {len(slide_images)} slides")
                return slide_images
            
            # Method 2: Direct LibreOffice PNG conversion (fallback)
            slide_images = await self._convert_via_libreoffice_png(pptx_path, project_dir)
            if slide_images:
                print(f"✅ LibreOffice PNG conversion successful: {len(slide_images)} slides")
                return slide_images
            
            # Method 3: pdf2image conversion (fallback)
            slide_images = await self._convert_via_pdf2image(pptx_path, project_dir)
            if slide_images:
                print(f"✅ pdf2image conversion successful: {len(slide_images)} slides")
                return slide_images
                
            # Final fallback: Create placeholder images
            print("⚠️ All conversion methods failed, creating placeholders")
            return await self._create_placeholder_images(pptx_path, project_dir)
                
        except Exception as e:
            print(f"❌ PPTX conversion error: {e}")
            return await self._create_placeholder_images(pptx_path, project_dir)
    
    async def _convert_via_libreoffice_pdf(self, pptx_path: Path, project_dir: Path) -> List[str]:
        """
        Method 1: LibreOffice + pdfplumber (highest quality)
        PPTX → PDF → high-quality images
        """
        try:
            import subprocess
            import pdfplumber
            from pdf2image import convert_from_path
            
            # Check if LibreOffice is available (try multiple possible commands)
            libreoffice_commands = ["soffice", "libreoffice", "/opt/homebrew/bin/soffice"]
            libreoffice_cmd = None
            
            for cmd in libreoffice_commands:
                try:
                    result = subprocess.run(
                        [cmd, "--version"], 
                        capture_output=True, 
                        text=True
                    )
                    if result.returncode == 0:
                        libreoffice_cmd = cmd
                        print(f"✅ Found LibreOffice: {cmd}")
                        break
                except FileNotFoundError:
                    continue
            
            if not libreoffice_cmd:
                print("⚠️ LibreOffice not found")
                return []
            
            slides_dir = project_dir / "slides"
            temp_pdf = project_dir / "temp_presentation.pdf"
            
            await self.send_progress("✅", f"Found LibreOffice: {libreoffice_cmd}")
            
            # Step 1: Convert PPTX to PDF using LibreOffice
            await self.send_progress("🔄", "Converting PPTX to PDF with LibreOffice...")
            subprocess.run([
                libreoffice_cmd, "--headless", "--convert-to", "pdf",
                "--outdir", str(project_dir), str(pptx_path)
            ], check=True, capture_output=True)
            
            # Find the generated PDF
            pdf_files = list(project_dir.glob("*.pdf"))
            if not pdf_files:
                await self.send_progress("❌", "PDF conversion failed - no PDF file generated")
                return []
            
            pdf_path = pdf_files[0]
            await self.send_progress("✅", "PDF generated")
            
            # Step 2: Convert PDF to high-quality images using pdf2image
            await self.send_progress("🔄", "Converting PDF to high-quality images...")
            images = convert_from_path(
                pdf_path,
                dpi=300,  # High DPI for quality
                fmt='PNG',
                thread_count=4,
                poppler_path=None  # Use system poppler
            )
            
            slide_images = []
            for i, image in enumerate(images):
                image_filename = f"slide_{i + 1:03d}.png"
                image_path = slides_dir / image_filename
                
                # Optimize and save image
                if image.mode != 'RGB':
                    image = image.convert('RGB')
                
                # Resize if too large (maintain aspect ratio)
                if image.width > 1920 or image.height > 1080:
                    image.thumbnail((1920, 1080), Image.Resampling.LANCZOS)
                
                image.save(image_path, "PNG", optimize=True, quality=95)
                slide_images.append(str(image_path))
                await self.send_progress("✅", f"Saved slide {i+1}")
            
            # Clean up temporary PDF
            if pdf_path.exists():
                pdf_path.unlink()
            
            return slide_images
            
        except subprocess.CalledProcessError as e:
            print(f"❌ LibreOffice conversion failed: {e}")
            return []
        except ImportError as e:
            print(f"❌ Missing dependency: {e}")
            return []
        except Exception as e:
            print(f"❌ LibreOffice + pdfplumber conversion failed: {e}")
            return []
    
    async def _convert_via_libreoffice_png(self, pptx_path: Path, project_dir: Path) -> List[str]:
        """
        Method 2: Direct LibreOffice PNG conversion (fallback)
        """
        try:
            import subprocess
            
            # Check if LibreOffice is available (try multiple possible commands)
            libreoffice_commands = ["soffice", "libreoffice", "/opt/homebrew/bin/soffice"]
            libreoffice_cmd = None
            
            for cmd in libreoffice_commands:
                try:
                    result = subprocess.run(
                        [cmd, "--version"], 
                        capture_output=True, 
                        text=True
                    )
                    if result.returncode == 0:
                        libreoffice_cmd = cmd
                        print(f"✅ Found LibreOffice: {cmd}")
                        break
                except FileNotFoundError:
                    continue
            
            if not libreoffice_cmd:
                print("⚠️ LibreOffice not found")
                return []
            
            slides_dir = project_dir / "slides"
            
            print("🔄 Converting PPTX to PNG with LibreOffice...")
            subprocess.run([
                libreoffice_cmd, "--headless", "--convert-to", "png",
                "--outdir", str(slides_dir), str(pptx_path)
            ], check=True, capture_output=True)
            
            # Find generated images and rename them properly
            png_files = sorted(list(slides_dir.glob("*.png")))
            slide_images = []
            
            for i, png_file in enumerate(png_files):
                new_filename = f"slide_{i + 1:03d}.png"
                new_path = slides_dir / new_filename
                
                if png_file != new_path:
                    png_file.rename(new_path)
                
                slide_images.append(str(new_path))
                print(f"✅ Processed slide {i+1}: {new_path}")
            
            return slide_images
            
        except subprocess.CalledProcessError as e:
            print(f"❌ LibreOffice PNG conversion failed: {e}")
            return []
        except Exception as e:
            print(f"❌ LibreOffice PNG conversion error: {e}")
            return []
    
    async def _convert_via_pdf2image(self, pptx_path: Path, project_dir: Path) -> List[str]:
        """
        Method 3: pdf2image conversion (fallback)
        First convert PPTX to PDF using python-pptx, then to images
        """
        try:
            from pdf2image import convert_from_path
            import tempfile
            
            # This is a basic fallback - would need additional logic
            # to convert PPTX to PDF first without LibreOffice
            print("⚠️ pdf2image fallback not fully implemented")
            return []
            
        except ImportError as e:
            print(f"❌ pdf2image not available: {e}")
            return []
        except Exception as e:
            print(f"❌ pdf2image conversion failed: {e}")
            return []

    async def _create_placeholder_images(self, pptx_path: Path, project_dir: Path) -> List[str]:
        """
        Create placeholder images when conversion fails
        """
        try:
            # Count slides in PPTX
            prs = Presentation(pptx_path)
            slide_count = len(prs.slides)
            
            slide_images = []
            slides_dir = project_dir / "slides"
            
            for i in range(slide_count):
                # Create a simple placeholder image
                img = Image.new('RGB', (1920, 1080), color='white')
                
                # You could add text or other elements here
                # For now, just save a white placeholder
                
                image_filename = f"slide_{i + 1:03d}.png"
                image_path = slides_dir / image_filename
                img.save(image_path, "PNG")
                
                slide_images.append(str(image_path))
            
            return slide_images
            
        except Exception as e:
            raise Exception(f"Placeholder image creation failed: {str(e)}")
    
    def get_project_metadata(self, project_id: str) -> Optional[dict]:
        """
        Get project metadata
        """
        try:
            metadata_path = self.upload_dir / project_id / "metadata.json"
            if metadata_path.exists():
                import json
                with open(metadata_path, 'r', encoding='utf-8') as f:
                    return json.load(f)
            return None
        except Exception:
            return None
    
    def cleanup_temp_files(self):
        """
        Clean up temporary files older than 1 hour
        """
        try:
            import time
            current_time = time.time()
            
            for file_path in self.temp_dir.iterdir():
                if file_path.is_file():
                    file_age = current_time - file_path.stat().st_mtime
                    if file_age > 3600:  # 1 hour
                        file_path.unlink()
                        
        except Exception:
            pass  # Ignore cleanup errors

# Global file service instance
file_service = FileService()
